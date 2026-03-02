"""
Generate PowerPoint slides for the Agentic AI Masterclass.
One deck per module. Reuses the design system from gen_slides.py.

Usage: python course/slides/gen_module_slides.py
Output: course/slides/Module_1_Foundations.pptx through Module_4_Capstones.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Design Constants ──────────────────────────────────

BG = RGBColor(10, 10, 46)
WHITE = RGBColor(255, 255, 255)
PURPLE = RGBColor(122, 95, 255)
CYAN = RGBColor(0, 210, 255)
PINK = RGBColor(255, 107, 157)
GRAY = RGBColor(170, 170, 170)
DARK_GRAY = RGBColor(100, 100, 100)
LIGHT = RGBColor(220, 220, 220)
GREEN = RGBColor(0, 255, 136)
ORANGE = RGBColor(255, 170, 0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Helpers ───────────────────────────────────────────


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, w, h, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return txBox


def add_bullets(slide, left, top, w, h, items, size=20, color=LIGHT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
    return txBox


# ─── Slide Templates ──────────────────────────────────


def title_slide(prs, title, subtitle="", section=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PURPLE)
    if section:
        add_text(slide, Inches(9), Inches(0.3), Inches(4), Inches(0.4),
                 section, size=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
             title, size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
                 subtitle, size=24, color=GRAY, align=PP_ALIGN.CENTER)
    return slide


def learning_slide(prs, title, items, section=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PURPLE)
    if section:
        add_text(slide, Inches(9), Inches(0.3), Inches(4), Inches(0.4),
                 section, size=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             title, size=36, color=CYAN, bold=True)
    add_bullets(slide, Inches(1.2), Inches(1.7), Inches(10.5), Inches(5), items)
    return slide


def lab_what_slide(prs, lab_name, prompt_text, what_agent_builds, section=""):
    """Per-lab: What To Do — prompt to paste, what agent builds."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PURPLE)
    if section:
        add_text(slide, Inches(9), Inches(0.3), Inches(4), Inches(0.4),
                 section, size=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             lab_name, size=36, color=PURPLE, bold=True)

    # Prompt box
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5),
              fill_color=RGBColor(20, 20, 60), border_color=CYAN)
    add_text(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.5),
             "Paste This Prompt:", size=16, color=CYAN, bold=True)
    add_text(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(3.5),
             prompt_text, size=14, color=LIGHT)

    # What gets built box
    add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.5),
              fill_color=RGBColor(20, 20, 60), border_color=GREEN)
    add_text(slide, Inches(7.0), Inches(1.7), Inches(5.4), Inches(0.5),
             "What the Agent Builds:", size=16, color=GREEN, bold=True)
    add_bullets(slide, Inches(7.0), Inches(2.3), Inches(5.4), Inches(3.5),
                what_agent_builds, size=14, color=LIGHT)
    return slide


def lab_run_slide(prs, lab_name, run_command, trace_output, section=""):
    """Per-lab: How to Run — terminal command + expected trace output."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PURPLE)
    if section:
        add_text(slide, Inches(9), Inches(0.3), Inches(4), Inches(0.4),
                 section, size=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             f"{lab_name} — How to Run", size=36, color=ORANGE, bold=True)

    # Command box
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
              fill_color=RGBColor(30, 30, 30), border_color=ORANGE)
    add_text(slide, Inches(1.0), Inches(1.6), Inches(11), Inches(0.6),
             f"$ {run_command}", size=20, color=GREEN, bold=True)

    # Trace output box
    add_shape(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(4.2),
              fill_color=RGBColor(15, 15, 35), border_color=GRAY)
    add_text(slide, Inches(1.0), Inches(2.7), Inches(11), Inches(0.4),
             "Expected Trace Output:", size=14, color=GRAY)
    add_text(slide, Inches(1.0), Inches(3.2), Inches(11), Inches(3.5),
             trace_output, size=12, color=LIGHT)
    return slide


def takeaways_slide(prs, lab_name, takeaways, section=""):
    """Per-lab: Key Takeaways."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PURPLE)
    if section:
        add_text(slide, Inches(9), Inches(0.3), Inches(4), Inches(0.4),
                 section, size=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             f"{lab_name} — Key Takeaways", size=36, color=GREEN, bold=True)
    add_bullets(slide, Inches(1.2), Inches(1.7), Inches(10.5), Inches(5),
                takeaways, size=22)
    return slide


def summary_slide(prs, title, rows, section=""):
    """Module summary table as a recap slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PURPLE)
    if section:
        add_text(slide, Inches(9), Inches(0.3), Inches(4), Inches(0.4),
                 section, size=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             title, size=36, color=PURPLE, bold=True)

    # Table as stacked cards
    y = Inches(1.6)
    for lab, pattern, insight in rows:
        add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.7),
                  fill_color=RGBColor(20, 20, 60), border_color=PURPLE)
        add_text(slide, Inches(1.0), y, Inches(2.5), Inches(0.7),
                 lab, size=16, color=CYAN, bold=True)
        add_text(slide, Inches(3.8), y, Inches(3.2), Inches(0.7),
                 pattern, size=16, color=LIGHT)
        add_text(slide, Inches(7.2), y, Inches(5), Inches(0.7),
                 insight, size=14, color=GRAY)
        y += Inches(0.85)
    return slide


# ═══════════════════════════════════════════════════════
# MODULE DATA
# ═══════════════════════════════════════════════════════

def gen_module_1(out_dir):
    prs = new_prs()
    sec = "Module 1: Foundations"

    title_slide(prs, "Module 1: Foundations",
                "From Shallow to Deep Agents", sec)

    learning_slide(prs, "What You'll Learn", [
        "The 4 pillars of deep agents: Planning, Sub-Agents, Filesystem, System Prompts",
        "create_deep_agent() — the factory function for building agents",
        "@tool decorator — how to give agents custom capabilities",
        "Thread IDs — conversation identity for stateful chats",
        "System prompts — shaping agent behavior with explicit instructions",
    ], sec)

    # Lab 0
    lab_what_slide(prs, "Lab 0: Environment Setup", (
        "The agent builds a setup verification\n"
        "script that checks Python version,\n"
        "installed packages, API key, and runs\n"
        "a smoke test agent."
    ), [
        "my-agents/lab0_setup.py",
        "Checks: Python 3.11+, deepagents, openai",
        "Smoke test: creates a test agent",
        "Prints ALL CHECKS PASSED",
    ], sec)

    lab_run_slide(prs, "Lab 0", "python my-agents/lab0_setup.py", (
        "[1] Python version: 3.13.x  ✓\n"
        "[2] deepagents package:  ✓\n"
        "[3] langchain-openai package:  ✓\n"
        "[4] OPENAI_API_KEY:  ✓ (starts with sk-...)\n"
        "[5] openai package:  ✓\n\n"
        "ALL CHECKS PASSED — Ready to go!\n\n"
        "[6] Testing agent creation...\n"
        "    Agent response: SETUP OK\n"
        "    ✓ Agent is working!"
    ), sec)

    # Lab 1
    lab_what_slide(prs, "Lab 1: Your First Agent", (
        "The agent builds a calculator + clock\n"
        "agent with two custom @tool functions.\n"
        "Also builds trace_utils.py for visible\n"
        "traces across all labs."
    ), [
        "my-agents/trace_utils.py (shared utility)",
        "my-agents/lab1_first_agent.py",
        "Tools: calculator(), get_current_time()",
        "4 test queries + interactive mode",
    ], sec)

    lab_run_slide(prs, "Lab 1", "python my-agents/lab1_first_agent.py", (
        "YOU: What is the square root of 1764?\n\n"
        "  [CALC] calculator(expression=\"sqrt(1764)\")\n"
        "  [RESULT] Result: 42.0\n\n"
        "AGENT RESPONSE:\n"
        "The square root of 1764 is 42.0\n\n"
        "══════════════════════════════════\n"
        "  YOUR TURN — Try it yourself!\n"
        "══════════════════════════════════\n"
        "> _"
    ), sec)

    takeaways_slide(prs, "Lab 1", [
        "@tool decorator turns Python functions into LLM-callable tools",
        "create_deep_agent() wires model + tools + system prompt together",
        "The LLM decides WHEN and HOW to call tools (you don't)",
        "System prompt shapes behavior: 'ALWAYS use calculator' prevents guessing",
    ], sec)

    summary_slide(prs, "Module 1 Summary", [
        ("Lab 0", "Setup Verifier", "Diagnostic script + smoke test"),
        ("Lab 1", "Custom Tools", "@tool, create_deep_agent(), thread_id"),
    ], sec)

    path = os.path.join(out_dir, "Module_1_Foundations.pptx")
    prs.save(path)
    print(f"  Saved: {path}")


def gen_module_2(out_dir):
    prs = new_prs()
    sec = "Module 2: Building Blocks"

    title_slide(prs, "Module 2: Building Blocks",
                "Tools, Planning, and Context Management", sec)

    learning_slide(prs, "What You'll Learn", [
        "Built-in filesystem tools: read_file, write_file, ls, glob, grep",
        "Planning with write_todos — structured task decomposition",
        "The write-then-summarize pattern — files survive context compression",
        "OpenAI web search via the Responses API",
        "System prompt workflows — numbered steps for reliable behavior",
    ], sec)

    lab_what_slide(prs, "Lab 2: Research Agent", (
        "The agent builds a web search +\n"
        "research report agent that plans,\n"
        "searches, takes notes to files, and\n"
        "writes a final report."
    ), [
        "my-agents/lab2_research_agent.py",
        "Custom tool: web_search() via OpenAI API",
        "5-step workflow in system prompt",
        "Files: /research/*.md, /output/report.md",
    ], sec)

    lab_run_slide(prs, "Lab 2", "python my-agents/lab2_research_agent.py", (
        "  [PLAN] write_todos([\"Search Python\", \"Search JS\", ...])\n"
        "  [SEARCH] web_search(query=\"Python backend 2025\")\n"
        "  [RESULT] Python continues to dominate...\n"
        "  [WRITE] write_file(\"/research/python.md\")\n"
        "  [FILE WRITTEN] /research/python.md\n"
        "  [SEARCH] web_search(query=\"JavaScript backend 2025\")\n"
        "  [RESULT] JavaScript with Node.js...\n"
        "  [WRITE] write_file(\"/output/report.md\")\n"
        "  [FILE WRITTEN] /output/report.md\n\n"
        "AGENT RESPONSE: Here is my comparison report..."
    ), sec)

    takeaways_slide(prs, "Lab 2", [
        "Filesystem tools are built-in — agents can read/write files automatically",
        "write_todos creates a plan before executing — prevents random-walk behavior",
        "Write immediately pattern: notes survive context summarization",
        "System prompt workflows (numbered steps) produce reliable, repeatable results",
    ], sec)

    summary_slide(prs, "Module 2 Summary", [
        ("Lab 2", "Research Agent", "Web search + filesystem + planning"),
    ], sec)

    path = os.path.join(out_dir, "Module_2_Building_Blocks.pptx")
    prs.save(path)
    print(f"  Saved: {path}")


def gen_module_3(out_dir):
    prs = new_prs()
    sec = "Module 3: Multi-Agent"

    title_slide(prs, "Module 3: Multi-Agent Orchestration",
                "Agents That Collaborate", sec)

    learning_slide(prs, "What You'll Learn", [
        "Sub-agent delegation with the task() tool",
        "Context isolation — sub-agents get fresh, separate contexts",
        "Parallel execution — multiple task() calls in one response",
        "Filesystem as a communication bus between agents",
        "Sequential pipelines — Agent A writes, Agent B reads",
    ], sec)

    # Lab 3a
    lab_what_slide(prs, "Lab 3a: Your First Sub-Agent", (
        "The agent builds an orchestrator\n"
        "with NO tools that MUST delegate\n"
        "all research to a sub-agent."
    ), [
        "my-agents/lab3a_first_subagent.py",
        "Orchestrator: gpt-4o, tools=[] (forced delegation)",
        "Researcher sub-agent: gpt-4o-mini, web_search",
        "Pattern: orchestrator → task() → researcher → answer",
    ], sec)

    lab_run_slide(prs, "Lab 3a", "python my-agents/lab3a_first_subagent.py", (
        '  [DELEGATE] task("Research LangGraph and LangChain",\n'
        '             subagent_type="researcher")\n'
        "  [SUB-AGENT DONE] LangGraph is a framework for building...\n\n"
        "AGENT RESPONSE:\n"
        "Based on the researcher's findings, LangGraph is..."
    ), sec)

    takeaways_slide(prs, "Lab 3a", [
        "tools=[] forces the orchestrator to delegate — it cannot do anything itself",
        "Sub-agents get isolated context (no cross-contamination)",
        "Only the final message returns (prevents context bloat)",
        "Different models per agent optimizes cost (gpt-4o orchestrator, gpt-4o-mini sub-agent)",
    ], sec)

    # Lab 3b
    lab_what_slide(prs, "Lab 3b: Parallel Sub-Agents", (
        "The agent builds an orchestrator\n"
        "that launches 3 specialist agents\n"
        "IN PARALLEL, then synthesizes."
    ), [
        "my-agents/lab3b_parallel_agents.py",
        "3 specialists: tech, market, pros/cons",
        "All 3 run simultaneously via parallel task() calls",
        "Orchestrator synthesizes into unified report",
    ], sec)

    lab_run_slide(prs, "Lab 3b", "python my-agents/lab3b_parallel_agents.py", (
        '  [DELEGATE] task("Technical analysis", subagent_type="tech_researcher")\n'
        '  [DELEGATE] task("Market analysis", subagent_type="market_researcher")\n'
        '  [DELEGATE] task("Pros/cons", subagent_type="pros_cons_analyst")\n'
        "  [SUB-AGENT DONE] Technical findings: Rust offers...\n"
        "  [SUB-AGENT DONE] Market data shows Go has...\n"
        "  [SUB-AGENT DONE] Strengths/weaknesses: Rust excels...\n\n"
        "AGENT RESPONSE: Synthesized report combining all 3 perspectives..."
    ), sec)

    takeaways_slide(prs, "Lab 3b", [
        "Multiple task() calls in ONE response = parallel execution",
        "3 specialists go deeper than 1 generalist (focused contexts)",
        "Orchestrator synthesizes multiple perspectives into unified report",
        "Each sub-agent writes to /research/ — orchestrator reads and combines",
    ], sec)

    # Lab 3c
    lab_what_slide(prs, "Lab 3c: Filesystem Communication", (
        "The agent builds a sequential pipeline\n"
        "where agents communicate ONLY through\n"
        "files — the filesystem is the message bus."
    ), [
        "my-agents/lab3c_file_sharing.py",
        "data_collector writes /data/*.md",
        "analyst reads /data/, calculates, writes /analysis/*.md",
        "Orchestrator reads /analysis/ and presents",
    ], sec)

    lab_run_slide(prs, "Lab 3c", "python my-agents/lab3c_file_sharing.py", (
        '  [DELEGATE] task("Collect pricing data", subagent_type="data_collector")\n'
        "  [SUB-AGENT DONE] Created /data/cloud_pricing.md\n"
        '  [DELEGATE] task("Analyze data", subagent_type="analyst")\n'
        "  [CALC] calculate(expression=\"67+34+25\")\n"
        "  [RESULT] 126\n"
        "  [SUB-AGENT DONE] Analysis written to /analysis/results.md\n\n"
        "AGENT RESPONSE: Based on the analysis, AWS is cheapest at..."
    ), sec)

    takeaways_slide(prs, "Lab 3c", [
        "The filesystem IS the message bus — agents communicate through files only",
        "Sequential pipeline: each agent depends on the previous agent's output",
        "Agents don't know about each other — only about file paths",
        "Data persists even after context summarization",
    ], sec)

    summary_slide(prs, "Module 3 Summary", [
        ("Lab 3a", "Delegation", "tools=[] forces task() delegation"),
        ("Lab 3b", "Parallel Execution", "Multiple task() calls = parallel"),
        ("Lab 3c", "File Pipeline", "Filesystem = async message bus"),
    ], sec)

    path = os.path.join(out_dir, "Module_3_Multi_Agent.pptx")
    prs.save(path)
    print(f"  Saved: {path}")


def gen_module_4(out_dir):
    prs = new_prs()
    sec = "Module 4: Capstones"

    title_slide(prs, "Module 4: Capstone Projects",
                "Production-Grade Multi-Agent Apps", sec)

    learning_slide(prs, "What You'll Learn", [
        "Multi-phase orchestration (parallel + sequential combined)",
        "Simulated data tools for educational environments",
        "SQLite integration for text-to-SQL agents",
        "Content pipeline with quality control (research -> write -> edit)",
        "Cost optimization: different models for different agent roles",
    ], sec)

    # Lab 4
    lab_what_slide(prs, "Lab 4: Financial Deep Research", (
        "The agent builds a 4-agent financial\n"
        "analysis system with phased execution:\n"
        "parallel data collection, then sequential\n"
        "analysis and report writing."
    ), [
        "my-agents/lab4_financial_research.py",
        "4 sub-agents: market_data, news, analyst, writer",
        "Phase 1 (parallel): data + news",
        "Phases 2-3 (sequential): analysis -> report",
        "Output: /report/investment_report.md",
    ], sec)

    lab_run_slide(prs, "Lab 4", "python my-agents/lab4_financial_research.py", (
        '  [DELEGATE] task("Fetch NVDA data", subagent_type="market_data_agent")\n'
        '  [DELEGATE] task("Research NVDA news", subagent_type="news_analyst")\n'
        "  [SUB-AGENT DONE] Market data collected for NVDA\n"
        "  [SUB-AGENT DONE] News analysis: Bullish sentiment...\n"
        '  [DELEGATE] task("Quantitative analysis", subagent_type="finance_analyst")\n'
        "  [SUB-AGENT DONE] Intrinsic value: $xxx, Risk: Medium\n"
        '  [DELEGATE] task("Write report", subagent_type="report_writer")\n'
        "  [SUB-AGENT DONE] Investment report written\n\n"
        "AGENT RESPONSE: Investment Analysis: NVIDIA (NVDA)..."
    ), sec)

    takeaways_slide(prs, "Lab 4", [
        "Multi-phase = parallel (Phase 1) + sequential (Phases 2-3) in one system",
        "4 specialist agents with different models optimize for cost vs. capability",
        "File-based communication scales: /data/market/, /data/news/, /data/analysis/, /report/",
        "Structured report template ensures professional, consistent output",
    ], sec)

    # Lab 5
    lab_what_slide(prs, "Lab 5: Text-to-SQL Analytics", (
        "The agent builds a conversational SQL\n"
        "agent with a sample e-commerce database.\n"
        "Natural language in -> SQL + results +\n"
        "explanations out."
    ), [
        "my-agents/lab5_text_to_sql.py",
        "SQLite database: products, customers, orders",
        "Tools: explore_schema(), run_sql_query()",
        "Safety: read-only, blocks write operations",
        "Persistent thread: agent remembers schema",
    ], sec)

    lab_run_slide(prs, "Lab 5", "python my-agents/lab5_text_to_sql.py", (
        "YOU: What are our top 5 products by total revenue?\n\n"
        "  [SCHEMA] explore_schema()\n"
        "  [RESULT] DATABASE SCHEMA: products (15 rows)...\n"
        '  [SQL] run_sql_query(query="SELECT p.name, SUM(p.price * o.quantity)...'
        '")\n'
        "  [RESULT] Mechanical Keyboard | $42,356.74 ...\n\n"
        "AGENT RESPONSE:\n"
        "The top 5 products by revenue are..."
    ), sec)

    takeaways_slide(prs, "Lab 5", [
        "Schema-first pattern: agent explores database before writing SQL",
        "Read-only safety: blocks DROP, DELETE, INSERT, UPDATE, ALTER, CREATE",
        "Persistent thread_id lets the agent remember schema across queries",
        "Natural language -> SQL -> Results -> Business explanation pipeline",
    ], sec)

    # Lab 6
    lab_what_slide(prs, "Lab 6: Content Pipeline", (
        "The agent builds a 3-stage content\n"
        "creation pipeline with quality control:\n"
        "Researcher -> Writer -> Editor."
    ), [
        "my-agents/lab6_content_pipeline.py",
        "3 sub-agents: researcher, writer, editor",
        "Sequential: each stage reads previous output",
        "Quality control: editor fact-checks against research",
        "Files: /research/ -> /drafts/ -> /review/ + /final/",
    ], sec)

    lab_run_slide(prs, "Lab 6", "python my-agents/lab6_content_pipeline.py", (
        '  [DELEGATE] task("Research AI agents", subagent_type="topic_researcher")\n'
        "  [SEARCH] web_search(\"AI agents software development 2025\")\n"
        "  [RESULT] AI agents are transforming...\n"
        "  [SUB-AGENT DONE] Research notes written to /research/\n"
        '  [DELEGATE] task("Write article", subagent_type="content_writer")\n'
        "  [SUB-AGENT DONE] Draft written to /drafts/article.md\n"
        '  [DELEGATE] task("Edit and review", subagent_type="editor_reviewer")\n'
        "  [SUB-AGENT DONE] Final article at /final/article.md\n\n"
        "AGENT RESPONSE: Here is the final article + editor's notes..."
    ), sec)

    takeaways_slide(prs, "Lab 6", [
        "Sequential pipeline ensures each stage has the data it needs",
        "Quality control: editor reviews draft against original research notes",
        "File-based handoff: /research/ -> /drafts/ -> /review/ + /final/",
        "Separation of concerns: research != writing != editing",
    ], sec)

    summary_slide(prs, "Module 4 Summary", [
        ("Lab 4", "Financial System", "4 agents, parallel + sequential phases"),
        ("Lab 5", "Text-to-SQL", "Schema exploration + safe SQL execution"),
        ("Lab 6", "Content Pipeline", "3-stage sequential + quality control"),
    ], sec)

    # Course completion slide
    title_slide(prs, "Course Complete!",
                "9 labs, 4 modules, 5 core patterns mastered", sec)

    path = os.path.join(out_dir, "Module_4_Capstones.pptx")
    prs.save(path)
    print(f"  Saved: {path}")


# ═══════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════

if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    print("Generating module slide decks...\n")
    gen_module_1(out_dir)
    gen_module_2(out_dir)
    gen_module_3(out_dir)
    gen_module_4(out_dir)
    print(f"\nDone! All 4 decks saved to {out_dir}/")
